"""
Simple Markdown-to-PowerPoint converter for slides separated by '---'.
Requires: python-pptx
Usage: python3 scripts/md_to_pptx.py
Generates: presentation_slides.pptx in repo root
"""
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path

IN = Path(__file__).resolve().parents[1] / "presentation_slides.md"
OUT = Path(__file__).resolve().parents[1] / "presentation_slides.pptx"

if not IN.exists():
    print(f"Input file not found: {IN}")
    raise SystemExit(1)

text = IN.read_text(encoding="utf-8")
# Split slides by lines that contain exactly '---'
slides_raw = []
current = []
for line in text.splitlines():
    if line.strip() == '---':
        if current:
            slides_raw.append("\n".join(current).strip())
            current = []
    else:
        current.append(line)
if current:
    slides_raw.append("\n".join(current).strip())

prs = Presentation()
# Use default slide width/height and layouts

for i, slide_md in enumerate(slides_raw):
    lines = [l.rstrip() for l in slide_md.splitlines() if l.strip() != '']
    if not lines:
        continue
    # Determine title
    title = None
    body_lines = []
    for ln in lines:
        if ln.startswith('# '):
            if title is None:
                title = ln.lstrip('# ').strip()
            else:
                body_lines.append(ln)
        elif ln.startswith('## '):
            if title is None:
                title = ln.lstrip('# ').strip()
            else:
                body_lines.append(ln)
        else:
            body_lines.append(ln)

    if i == 0 and title:
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        # subtitle if present: take first non-title line
        subtitle = None
        for bl in body_lines:
            if bl.strip():
                subtitle = bl
                break
        if subtitle and slide.placeholders and len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = subtitle
            except Exception:
                pass
        continue

    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    if title:
        slide.shapes.title.text = title
    else:
        # fallback: first body line as title
        slide.shapes.title.text = body_lines[0] if body_lines else ''
        body_lines = body_lines[1:]

    # Find content placeholder
    body = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 2:
            body = shape
            break
    if body is None:
        # fallback: add a textbox
        from pptx.util import Inches
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        body = slide.shapes.add_textbox(left, top, width, height)

    tf = body.text_frame
    tf.clear()

    # Parse body lines into bullets / codeblocks
    in_code = False
    for bl in body_lines:
        if bl.strip().startswith('```'):
            in_code = not in_code
            if in_code:
                p = tf.add_paragraph()
                p.level = 0
                p.font.name = 'Courier New'
                p.font.size = Pt(12)
                p.text = ''
            continue
        if in_code:
            # append code as new paragraph with monospace
            p = tf.add_paragraph()
            p.level = 0
            p.font.name = 'Courier New'
            p.font.size = Pt(12)
            p.text = bl
            continue

        stripped = bl.lstrip()
        if stripped.startswith('- '):
            p = tf.add_paragraph()
            p.level = 0
            p.font.size = Pt(18)
            p.text = stripped[2:].strip()
        elif stripped.startswith('* '):
            p = tf.add_paragraph()
            p.level = 1
            p.font.size = Pt(16)
            p.text = stripped[2:].strip()
        elif stripped.startswith('## '):
            p = tf.add_paragraph()
            p.level = 0
            p.font.size = Pt(20)
            p.text = stripped.lstrip('# ').strip()
        else:
            # plain paragraph
            p = tf.add_paragraph()
            p.level = 0
            p.font.size = Pt(14)
            p.text = bl

prs.save(str(OUT))
print(f"Saved: {OUT}")
